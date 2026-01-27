import streamlit as st
import os
import json
from google import genai
from google.genai import types
from docxtpl import DocxTemplate, RichText
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import re
from datetime import datetime

# --- CONFIGURAZIONE PAGINA ---
st.set_page_config(page_title="Generatore Documenti AI", page_icon="🪄", layout="wide")
st.title("🪄 Generatore Report GEOCOM® (Word & PPT)")

# --- CONFIGURAZIONE API KEY ---
API_KEY_FISSA = "" 
api_key = API_KEY_FISSA if API_KEY_FISSA else st.sidebar.text_input("Google AI Studio Key", type="password")

# --- SELETTORE FORMATO ---
formato_scelto = st.radio("Scegli il formato:", ["Documento Word (.docx)", "Presentazione PowerPoint (.pptx)"])

# --- INPUT UTENTE ---
testo_grezzo = st.text_area("Incolla qui il testo:", height=350)

# --- FUNZIONI DI SUPPORTO ---

def get_gemini_data(text, key, formato):
    client = genai.Client(api_key=key)
    if formato == "ppt":
        # PROMPT POTENZIATO PER PULIRE IL TESTO
        prompt_finale = f"""
        Sei un formattatore di dati per PowerPoint.
        IL TUO COMPITO:
        1. Analizza il testo in input. Spesso i punti elenco sono sulla stessa riga (es: "* A * B"). DEVI SEPARARLI con "\\n".
        2. RIMUOVI gli asterischi iniziali (*) dai punti elenco nel campo 'testo', perché PowerPoint mette già i pallini.
        3. Mantieni **grassetto** e *corsivo* all'interno delle frasi.
        
        Restituisci ESCLUSIVAMENTE un oggetto JSON (no liste, no markdown ```json):
        {{
            "titolo_report": "Titolo principale",
            "sottotitolo_report": "Sottotitolo eventuale",
            "lista_sezioni": [ 
                {{ 
                    "titolo": "Titolo della Slide (es. Risultati Ottenuti)", 
                    "testo": "Frase 1\\nFrase 2\\nFrase 3" 
                }} 
            ]
        }}
        TESTO DA ELABORARE:
        {text}
        """
    else:
        prompt_finale = f"Estrai dati per Word in JSON (usa ** per grassetti): {text}"

    response = client.models.generate_content(
        model='gemini-flash-latest', 
        contents=prompt_finale,
        config=types.GenerateContentConfig(response_mime_type='application/json')
    )
    
    try:
        res_data = json.loads(response.text)
        if isinstance(res_data, list):
            res_data = res_data[0]
        return res_data
    except Exception as e:
        st.error(f"Errore interpretazione AI: {e}. Riprova.")
        return {}

def add_formatted_text(paragraph, text):
    """Gestisce grassetto e corsivo dentro PowerPoint"""
    # Rimuove eventuali asterischi a inizio riga rimasti per errore
    text = text.lstrip('*').strip()
    
    # Regex per trovare **bold**, *italic* e ***bolditalic***
    pattern = re.compile(r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*)')
    parts = pattern.split(text)
    
    for part in parts:
        if not part: continue
        run = paragraph.add_run()
        if part.startswith('***') and part.endswith('***'):
            run.text = part[3:-3]
            run.font.bold = True
            run.font.italic = True
        elif part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith('*') and part.endswith('*'):
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part

def generate_doc(data):
    if not os.path.exists("template_aziendale.docx"): return None
    doc = DocxTemplate("template_aziendale.docx")
    if 'lista_sezioni' in data:
        for s in data['lista_sezioni']:
            rt = RichText()
            parts = re.split(r'(\*\*.*?\*\*)', s.get('testo', ''))
            for p in parts:
                if p.startswith('**'): rt.add(p[2:-2], bold=True)
                else: rt.add(p)
            s['testo'] = rt
    doc.render(data)
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

def generate_ppt(data):
    if not os.path.exists("template_aziendale.pptx"): return None
    prs = Presentation("template_aziendale.pptx")
    testo_fisso = "GEOCOM® - Report Analitico"

    # --- 1. SLIDE COPERTINA ---
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title: 
            slide.shapes.title.text = data.get('titolo_report', 'Report')
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"{data.get('sottotitolo_report', '')}\n{data.get('data_odierna', '')}"
    except: pass

    # --- 2. SLIDE CONTENUTO ---
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    for sezione in data.get('lista_sezioni', []):
        slide = prs.slides.add_slide(layout)
        
        # Titolo Slide
        if slide.shapes.title:
            slide.shapes.title.text = sezione.get('titolo', '')

        # --- GESTIONE TESTO FISSO ---
        # Creiamo una Textbox manuale in alto
        tb_fisso = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
        p_fisso = tb_fisso.text_frame.paragraphs[0]
        p_fisso.text = testo_fisso
        p_fisso.font.bold = True
        p_fisso.font.size = Pt(12)
        # Ho rimosso l'assegnazione errata del colore. 
        # Se vuoi forzare il nero, decommenta la riga sotto:
        # p_fisso.font.color.rgb = RGBColor(0, 0, 0)

        # --- GESTIONE CONTENUTO ---
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            testo_raw = sezione.get('testo', '')
            lines = testo_raw.split('\n')
            
            for line in lines:
                if line.strip():
                    p = tf.add_paragraph()
                    add_formatted_text(p, line)
                    p.level = 0 

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

# --- INTERFACCIA UTENTE ---
st.markdown("---") 

if st.button("🚀 Genera Documento", type="primary"):
    if not api_key:
        st.error("⚠️ Inserisci la API Key!")
    elif not testo_grezzo:
        st.warning("⚠️ Inserisci il testo da analizzare!")
    else:
        tipo = "ppt" if "PowerPoint" in formato_scelto else "word"
        
        with st.spinner("L'AI sta strutturando il report..."):
            try:
                dati = get_gemini_data(testo_grezzo, api_key, tipo)
                dati['data_odierna'] = datetime.now().strftime("%d/%m/%Y")
                
                if tipo == "word":
                    output = generate_doc(dati)
                    nome_file = "Report_Geocom.docx"
                    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                else:
                    output = generate_ppt(dati)
                    nome_file = "Report_Geocom.pptx"
                    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                
                if output:
                    st.success("✅ Generazione completata!")
                    st.download_button(
                        label="📥 Scarica Documento",
                        data=output,
                        file_name=nome_file,
                        mime=mime_type
                    )
            except Exception as e:
                st.error(f"❌ Si è verificato un errore: {e}")
                ## per runnarlo ---> python -m streamlit run formattazione4.py